from app.models.content import Content
from pptx.exc import PackageNotFoundError
from .handler import FileHandler
from typing import Optional, Any, Dict
from pptx import Presentation


class PPTXHandler(FileHandler):
    def __init__(self, filename: str):
        super().__init__(filename)

    def extract_metadata(self) -> Optional[dict[str, Any]]:
        # TODO: figure out why the metadata came empty in the test files.
        try:
            prop = Presentation(str(self.file)).core_properties
            metadata = {}
            for d in dir(prop):
                if not d.startswith('_'):
                    metadata[d] = getattr(prop, d)
            return metadata
        except PackageNotFoundError:
            raise FileNotFoundError
        except Exception as e:
            self.logger.warning(e)

    def extract_content(self) -> Content:
        def _extract_text(presentation) -> str:
            texts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)

        def _extract_tables(presentation) -> list:
            tables = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 19:  # 19 corresponds to the table shape type
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        tables.append(table_data)
            return tables

        presentation = Presentation(self.file)
        text_content = _extract_text(presentation)
        table_content = _extract_tables(presentation)
        return Content(text_content, table_content)
